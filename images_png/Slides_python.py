import os
from pptx import Presentation
from pptx.util import Inches
import re

def natural_sort_key(s):
    # Extrait les parties numériques pour un tri naturel
    return [int(text) if text.isdigit() else text.lower() for text in re.split('(\d+)', s)]

def create_powerpoint_from_png(png_folder, output_pptx):
    png_files = [f for f in os.listdir(png_folder) if f.endswith('.png')]
    
    # Trier les fichiers PNG avec une fonction de tri naturelle
    png_files.sort(key=natural_sort_key)

    # Créer une présentation PowerPoint
    presentation = Presentation()

    for png_file in png_files:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Utiliser une mise en page vide
        img_path = os.path.join(png_folder, png_file)
        
        # Ajouter l'image à la diapositive
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=presentation.slide_width, height=presentation.slide_height)
        
        # Renommer la diapositive (facultatif)
        step_index = png_file.split('_')[-1].split('.')[0]  # Extrait l'index de l'étape
        slide_name = f"carte_vents_step_{int(step_index):02d}"
        slide.name = slide_name  # Renomme la diapositive

    # Enregistrer la présentation
    presentation.save(output_pptx)

# Utilisation de la fonction
png_folder = r"C:\Users\arthu\OneDrive\Arthur\Programmation\TIPE_Arthur_Lhoste\images_png"
output_pptx = r"C:\Users\arthu\OneDrive\Arthur\Programmation\TIPE_Arthur_Lhoste\images_png\carte_vents.pptx"

png_folder1 = r"C:\Users\arthu\OneDrive\Arthur\Programmation\TIPE_Arthur_Lhoste\route_ideale"
output_pptx1 = r"C:\Users\arthu\OneDrive\Arthur\Programmation\TIPE_Arthur_Lhoste\route_ideale\carte_route_ideale.pptx"

create_powerpoint_from_png(png_folder1, output_pptx1)
